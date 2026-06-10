"""
Create PowerPoint presentations for all 4 quizzes
Format: Slide 1 = Question, Slide 2 = Answer, repeat...
"""

import sys
if sys.platform == 'win32':
    sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_quiz_ppt(quiz_data, filename, quiz_title):
    """Create a PowerPoint for a quiz with question/answer alternating slides"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(title_slide_layout)

    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = quiz_title
    title_frame.paragraphs[0].font.size = Pt(54)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 126)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"{len(quiz_data)} Questions"
    subtitle_frame.paragraphs[0].font.size = Pt(32)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add each question and answer as alternating slides
    for i, q in enumerate(quiz_data, 1):
        # Question slide
        q_slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background color for question
        background = q_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(232, 245, 233)  # Light green

        # Question number
        q_num_box = q_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.6))
        q_num_frame = q_num_box.text_frame
        q_num_frame.text = f"Question {i}"
        q_num_frame.paragraphs[0].font.size = Pt(28)
        q_num_frame.paragraphs[0].font.bold = True
        q_num_frame.paragraphs[0].font.color.rgb = RGBColor(27, 94, 32)

        # Question text
        q_text_box = q_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
        q_text_frame = q_text_box.text_frame
        q_text_frame.text = q['question']
        q_text_frame.paragraphs[0].font.size = Pt(36)
        q_text_frame.paragraphs[0].font.bold = True
        q_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        q_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        q_text_frame.word_wrap = True

        # Options (if multiple choice)
        if 'options' in q and q['options']:
            options_box = q_slide.shapes.add_textbox(Inches(1.5), Inches(5.2), Inches(7), Inches(2))
            options_frame = options_box.text_frame
            for j, opt in enumerate(q['options'], 1):
                p = options_frame.add_paragraph() if j > 1 else options_frame.paragraphs[0]
                p.text = f"{chr(64+j)}) {opt}"
                p.font.size = Pt(22)
                p.space_after = Pt(8)

        # Answer slide
        a_slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background color for answer
        background = a_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(227, 242, 253)  # Light blue

        # "Answer" label
        a_label_box = a_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(2), Inches(0.6))
        a_label_frame = a_label_box.text_frame
        a_label_frame.text = f"Answer {i}"
        a_label_frame.paragraphs[0].font.size = Pt(28)
        a_label_frame.paragraphs[0].font.bold = True
        a_label_frame.paragraphs[0].font.color.rgb = RGBColor(13, 71, 161)

        # Correct answer
        a_text_box = a_slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        a_text_frame = a_text_box.text_frame
        a_text_frame.text = f"✓ {q['answer']}"
        a_text_frame.paragraphs[0].font.size = Pt(40)
        a_text_frame.paragraphs[0].font.bold = True
        a_text_frame.paragraphs[0].font.color.rgb = RGBColor(27, 94, 32)
        a_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        a_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        a_text_frame.word_wrap = True

        # Explanation (if available)
        if 'explanation' in q and q['explanation']:
            exp_box = a_slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(2))
            exp_frame = exp_box.text_frame
            exp_frame.text = q['explanation']
            exp_frame.paragraphs[0].font.size = Pt(18)
            exp_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)
            exp_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            exp_frame.word_wrap = True

    # Save
    prs.save(filename)
    print(f"[+] Created: {filename} ({len(quiz_data)} Q&A pairs = {len(prs.slides)} slides)")

# Quiz 1: RAG Basics
quiz1_data = [
    {
        'question': 'What does RAG stand for?',
        'options': ['Retrieval Augmented Generation', 'Rapid AI Generation', 'Random Access Generation', 'Recursive Algorithm Generator'],
        'answer': 'Retrieval Augmented Generation',
        'explanation': 'RAG = Retrieval (search) + Augmented (add context) + Generation (LLM answers)'
    },
    {
        'question': 'Which problem does RAG primarily solve?',
        'options': ['Reducing hallucinations and providing up-to-date info', 'Making LLMs run faster', 'Making LLMs smaller', 'Training LLMs from scratch'],
        'answer': 'Reducing hallucinations and providing up-to-date information',
        'explanation': 'RAG grounds LLM answers in retrieved documents, reducing made-up information'
    },
    {
        'question': 'RAG has two main phases. What are they?',
        'options': ['Indexing and Query', 'Training and Testing', 'Reading and Writing', 'Input and Output'],
        'answer': 'Indexing and Query',
        'explanation': 'Indexing = prepare documents. Query = answer questions using those documents.'
    },
    {
        'question': 'In the Indexing phase, what do we NOT do?',
        'options': ['Answer user questions', 'Split documents into chunks', 'Create embeddings', 'Store in vector database'],
        'answer': 'Answer user questions',
        'explanation': 'Answering happens in Query phase! Indexing just prepares the knowledge base.'
    },
    {
        'question': 'True or False: RAG requires retraining the LLM on new data',
        'options': ['False', 'True'],
        'answer': 'False',
        'explanation': 'That\'s the beauty of RAG! No retraining needed. Just update your documents.'
    },
    {
        'question': 'Which analogy best describes RAG?',
        'options': ['Open book exam', 'Closed book exam', 'Pop quiz', 'Group project'],
        'answer': 'Open book exam (can reference materials)',
        'explanation': 'LLM can "reference" documents to find answers, like an open book exam!'
    },
    {
        'question': 'What is the role of the vector database in RAG?',
        'options': ['Store and search embeddings', 'Store LLM weights', 'Train embedding model', 'Generate final answer'],
        'answer': 'Store and search embeddings for relevant documents',
        'explanation': 'Vector DB enables fast semantic similarity search of document chunks'
    },
    {
        'question': 'In the paper exercise, Sheet B was easier because:',
        'options': ['Info organized with metadata', 'Fewer words', 'Larger font', 'Different questions'],
        'answer': 'Information was organized with metadata (like RAG!)',
        'explanation': 'Organization + metadata = fast retrieval. That\'s how RAG works!'
    }
]

# Quiz 2: Chunking
quiz2_data = [
    {
        'question': 'Why is chunking necessary for RAG?',
        'options': ['LLMs can\'t read entire docs at once', 'Makes docs prettier', 'Reduces storage', 'Speeds up internet'],
        'answer': 'LLMs cannot read entire documents at once',
        'explanation': 'LLMs have context limits. We chunk to fit within those limits and retrieve precisely.'
    },
    {
        'question': 'What happens if chunks are TOO SMALL?',
        'options': ['Lose context and meaning', 'Nothing wrong', 'Vector DB crashes', 'LLM generates longer answers'],
        'answer': 'Lose context and meaning',
        'explanation': 'Small chunks lack context. "The price is $50" without knowing what product!'
    },
    {
        'question': 'What happens if chunks are TOO LARGE?',
        'options': ['Irrelevant info included, confuses LLM', 'Better results always', 'Faster retrieval', 'More accurate embeddings'],
        'answer': 'Irrelevant information gets included, confuses the LLM',
        'explanation': 'Large chunks add noise. Too much info = harder for LLM to focus on the answer.'
    },
    {
        'question': 'Recommended chunk size for general text:',
        'options': ['200-500 tokens', '10-20 tokens', '50-100 tokens', '5000+ tokens'],
        'answer': '200-500 tokens',
        'explanation': 'Sweet spot! Enough context but not too much noise. ~1-2 paragraphs.'
    },
    {
        'question': 'What is chunk overlap and why is it useful?',
        'options': ['Duplicate text between chunks for context', 'Error in chunking', 'Compression technique', 'Method to save storage'],
        'answer': 'Duplicate text between chunks to preserve context across boundaries',
        'explanation': 'Overlap ensures info at chunk boundaries isn\'t lost. 10-20% overlap is typical.'
    },
    {
        'question': 'Which chunking strategy is BEST for most scenarios?',
        'options': ['Recursive chunking', 'Fixed-size', 'Random', 'No chunking'],
        'answer': 'Recursive chunking',
        'explanation': 'Smart fallbacks! Tries paragraphs → sentences → words. Respects structure.'
    },
    {
        'question': 'Legal documents where accuracy is CRITICAL - what to prioritize?',
        'options': ['Precision - semantic chunking + metadata', 'Speed - largest chunks', 'Cost - smallest chunks', 'Simplicity - no chunking'],
        'answer': 'Precision - use semantic chunking + metadata for attribution',
        'explanation': 'Legal needs accuracy! Semantic chunking + metadata shows exact source/page.'
    },
    {
        'question': 'What is metadata in the context of RAG chunks?',
        'options': ['Data about data (source, page, date)', 'Secret information', 'Embedding dimensions', 'LLM temperature'],
        'answer': 'Data about data (e.g., source, page number, date)',
        'explanation': 'Metadata = extra info about chunk. Helps with filtering and attribution.'
    },
    {
        'question': 'For code chunking, best practice is to:',
        'options': ['Split by function/class boundaries', 'Split every 100 chars', 'Never split code', 'Split randomly'],
        'answer': 'Split by function/class boundaries',
        'explanation': 'Respect code structure! Keep functions intact for context.'
    },
    {
        'question': 'Recommended chunk overlap percentage:',
        'options': ['10-20%', '0%', '50%', '90%'],
        'answer': '10-20%',
        'explanation': 'Just enough overlap to preserve context at boundaries. Not too much waste.'
    }
]

# Quiz 3: Vector DB & Retrieval
quiz3_data = [
    {
        'question': 'What are embeddings in RAG?',
        'options': ['Numerical representations capturing semantic meaning', 'Compressed documents', 'Index numbers', 'File sizes'],
        'answer': 'Numerical representations of text that capture semantic meaning',
        'explanation': 'Embeddings = text converted to vectors. Similar meanings → similar vectors!'
    },
    {
        'question': 'Vector databases enable what type of search?',
        'options': ['Semantic similarity search', 'Exact keyword match only', 'Alphabetical sorting', 'Random selection'],
        'answer': 'Semantic similarity search',
        'explanation': 'Find similar MEANING, not just exact keywords. "Happy" finds "joyful"!'
    },
    {
        'question': 'Which is a popular vector database?',
        'options': ['ChromaDB', 'MySQL', 'Excel', 'Notepad'],
        'answer': 'ChromaDB',
        'explanation': 'ChromaDB, FAISS, Pinecone, Weaviate are all vector databases. Not regular DBs!'
    },
    {
        'question': 'Cosine similarity measures:',
        'options': ['Angle between vectors (semantic similarity)', 'File size difference', 'Number of words', 'Document age'],
        'answer': 'Angle between vectors (semantic similarity)',
        'explanation': 'Cosine similarity = angle between vectors. Smaller angle = more similar meaning!'
    },
    {
        'question': 'What does "k" in top-k retrieval mean?',
        'options': ['Number of most similar chunks to retrieve', 'Kilobytes of data', 'Thousand documents', 'Keyword count'],
        'answer': 'Number of most similar chunks to retrieve',
        'explanation': 'k=3 means retrieve top 3 most similar chunks. Typical range: 2-5.'
    },
    {
        'question': 'MMR (Maximal Marginal Relevance) is used for:',
        'options': ['Getting diverse results (not all similar)', 'Faster search', 'Smaller file sizes', 'Exact matching'],
        'answer': 'Getting diverse results (not all similar)',
        'explanation': 'MMR balances relevance + diversity. Avoids 5 chunks all saying same thing!'
    },
    {
        'question': 'Metadata filtering in vector search allows you to:',
        'options': ['Search only within specific categories', 'Delete unwanted results', 'Change vector dimensions', 'Increase speed always'],
        'answer': 'Search only within specific categories/sources',
        'explanation': 'Filter before searching! "Only search product docs, not customer stories"'
    },
    {
        'question': 'ChromaDB vs FAISS - main difference?',
        'options': ['ChromaDB easier, FAISS faster for production', 'No difference', 'FAISS always better', 'ChromaDB Windows only'],
        'answer': 'ChromaDB easier for prototyping, FAISS faster for production',
        'explanation': 'ChromaDB = simple to start. FAISS = optimized for billions of vectors.'
    },
    {
        'question': 'Why use similarity threshold in retrieval?',
        'options': ['Filter out low-quality matches', 'Speed up search', 'Reduce database size', 'Increase results'],
        'answer': 'Filter out low-quality matches',
        'explanation': 'Only return results above threshold. Better to say "don\'t know" than give bad answer!'
    },
    {
        'question': 'Hybrid search combines:',
        'options': ['Vector similarity + keyword matching', 'Two different LLMs', 'Multiple databases', 'SQL and NoSQL'],
        'answer': 'Vector similarity + keyword matching',
        'explanation': 'Best of both! Semantic meaning + exact term matching = powerful search.'
    }
]

# Quiz 4: Advanced & Evaluation
quiz4_data = [
    {
        'question': 'Precision in RAG evaluation measures:',
        'options': ['Of retrieved docs, how many are relevant?', 'Total number of docs', 'Speed of retrieval', 'Size of answers'],
        'answer': 'Of retrieved docs, how many are relevant?',
        'explanation': 'Precision = relevant_retrieved / total_retrieved. Quality over quantity!'
    },
    {
        'question': 'Recall in RAG evaluation measures:',
        'options': ['Of all relevant docs, how many were retrieved?', 'How fast we retrieve', 'Total DB size', 'Answer length'],
        'answer': 'Of all relevant docs, how many were retrieved?',
        'explanation': 'Recall = relevant_retrieved / total_relevant. Did we find everything important?'
    },
    {
        'question': 'Faithfulness metric checks if:',
        'options': ['Answer based on retrieved docs (no hallucination)', 'Answer is long enough', 'Retrieval was fast', 'User is satisfied'],
        'answer': 'Answer is based on retrieved documents (no hallucination)',
        'explanation': 'Faithfulness = no making things up! Answer must come from retrieved docs.'
    },
    {
        'question': 'RAGAS framework is used for:',
        'options': ['Automated RAG evaluation', 'Building vector databases', 'Training LLMs', 'Chunking documents'],
        'answer': 'Automated RAG evaluation',
        'explanation': 'RAGAS = RAG Assessment. Automates measuring retrieval + generation quality.'
    },
    {
        'question': 'Production RAG systems MUST have:',
        'options': ['Logging and monitoring', 'Fancy UI only', 'No error handling', 'Largest chunks'],
        'answer': 'Logging and monitoring',
        'explanation': 'Track performance, quality, costs! Can\'t improve what you don\'t measure.'
    },
    {
        'question': 'Caching in RAG helps with:',
        'options': ['Reducing costs and improving speed', 'Making answers longer', 'Increasing hallucinations', 'Deleting old data'],
        'answer': 'Reducing costs and improving speed',
        'explanation': 'Cache embeddings & frequent queries. Don\'t recompute same things!'
    },
    {
        'question': 'User feedback in RAG is important for:',
        'options': ['Improving system quality over time', 'Just collecting ratings', 'Slowing down responses', 'Making it complex'],
        'answer': 'Improving system quality over time',
        'explanation': 'Thumbs up/down → identify problems → improve chunking/prompts → better system!'
    },
    {
        'question': 'Answer Relevance metric checks if:',
        'options': ['Answer actually addresses the question', 'Answer is grammatically correct', 'Answer is short', 'Retrieval was accurate'],
        'answer': 'Answer actually addresses the question asked',
        'explanation': 'Did we answer the actual question? Or ramble about something else?'
    },
    {
        'question': 'Context Relevance metric checks if:',
        'options': ['Retrieved chunks are useful for answering', 'LLM is fast', 'Database is large', 'User is happy'],
        'answer': 'Retrieved chunks are actually useful for answering',
        'explanation': 'Were retrieved docs helpful? Or just noise that confused the LLM?'
    },
    {
        'question': 'MRR (Mean Reciprocal Rank) measures:',
        'options': ['How quickly we find first relevant result', 'Total number of results', 'Database size', 'Answer quality'],
        'answer': 'How quickly we find the first relevant result',
        'explanation': 'MRR = 1/rank of first good result. Rank 1 = perfect! Rank 5 = not great.'
    },
    {
        'question': 'Best practice for production RAG:',
        'options': ['All of the above', 'Log every query', 'Monitor metrics', 'Collect user feedback'],
        'answer': 'All of the above',
        'explanation': 'Production = logging + monitoring + feedback + error handling + caching!'
    },
    {
        'question': 'Temperature in LLM controls:',
        'options': ['Randomness (low=focused, high=creative)', 'Speed of generation', 'Number of tokens', 'Database temp'],
        'answer': 'Randomness (low=focused, high=creative)',
        'explanation': 'Temperature 0.0 = deterministic. Temperature 1.0 = very creative/random.'
    }
]

# Create all 4 PPTs
print("\n[*] Creating PowerPoint presentations for all quizzes...")
print("-" * 70)

create_quiz_ppt(quiz1_data, "Quiz_1_RAG_Basics.pptx", "Quiz 1: RAG Basics")
create_quiz_ppt(quiz2_data, "Quiz_2_Chunking.pptx", "Quiz 2: Chunking Strategies")
create_quiz_ppt(quiz3_data, "Quiz_3_VectorDB.pptx", "Quiz 3: Vector DB & Retrieval")
create_quiz_ppt(quiz4_data, "Quiz_4_Advanced.pptx", "Quiz 4: Evaluation & Production")

print("-" * 70)
print("[+] Done! All 4 PowerPoint presentations created.")
print("\nSlide format: Question (green) → Answer (blue) → repeat")
print("Total slides: Quiz1=17, Quiz2=21, Quiz3=21, Quiz4=25")
print("\nReady for your workshop presentations!")
